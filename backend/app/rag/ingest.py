from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.core import config


def load_pdf(path: str) -> str:
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
}


def load_document(path: str) -> str:
    ext = Path(path).suffix.lower()
    loader = LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported document type: {ext}")
    text = loader(path)
    if not text.strip():
        raise ValueError(f"No extractable text found in {path}")
    return text


def chunk_text(text: str, chunk_size_words: int = None, overlap_words: int = None) -> list[str]:
    chunk_size_words = chunk_size_words or config.CHUNK_SIZE_WORDS
    overlap_words = overlap_words or config.CHUNK_OVERLAP_WORDS
    words = text.split()

    chunks = []
    start = 0
    step = max(chunk_size_words - overlap_words, 1)
    while start < len(words):
        chunk = " ".join(words[start : start + chunk_size_words])
        if chunk.strip():
            chunks.append(chunk)
        start += step
    return chunks
